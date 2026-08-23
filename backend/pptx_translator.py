from pptx import Presentation
from langchain_openai import ChatOpenAI
import json
import os
import concurrent.futures

def _get_llm(model_name: str = "openai/gpt-5.6-luna"):
    return ChatOpenAI(
        model=model_name,
        api_key=os.getenv("OPENAI_API_KEY"),
        base_url=os.getenv("OPENAI_BASE_URL"),
        temperature=0.1
    )

def translate_text_batch(texts: list[str], target_language: str, model_name: str) -> list[str]:
    if not texts:
        return []
        
    llm = _get_llm(model_name)
    
    # We pass a JSON array of strings to the LLM and expect a JSON array back
    prompt = f"""
    You are an expert translator. Translate the following JSON array of text strings into {target_language}.
    Maintain the exact tone, intent, and whitespace where possible. 
    Return ONLY a valid JSON array of the translated strings in the exact same order.
    Do not add markdown formatting or explanation.

    Original JSON:
    {json.dumps(texts)}
    """
    
    try:
        response = llm.invoke(prompt)
        text = response.content.strip()
        if text.startswith("```json"):
            text = text[7:-3].strip()
        elif text.startswith("```"):
            text = text[3:-3].strip()
            
        translated = json.loads(text)
        
        # Fallback if the array length doesn't match
        if len(translated) != len(texts):
            print(f"Warning: Translated array length ({len(translated)}) doesn't match original ({len(texts)})")
            # Fall back to returning original text if completely broken
            return texts
            
        return translated
    except Exception as e:
        print(f"Translation batch failed: {e}")
        raise Exception("Translation API is currently unavailable or returned an error.")

def extract_and_translate_pptx(input_path: str, output_path: str, target_language: str, model_name: str = "openai/gpt-5.6-luna"):
    """
    Translates a PowerPoint presentation into a target language while preserving run-level formatting
    by capturing the format of the first run of a paragraph and applying it to the translated text.
    """
    prs = Presentation(input_path)
    
    # Collect all paragraphs that have text
    paragraphs_to_translate = []
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text.strip():
                        paragraphs_to_translate.append(paragraph)
            elif shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            if paragraph.text.strip():
                                paragraphs_to_translate.append(paragraph)
    
    # Extract strings
    original_texts = [p.text for p in paragraphs_to_translate]
    
    # Translate in batches of 20 to avoid exceeding context or breaking JSON arrays
    batch_size = 20
    translated_texts = []
    
    with concurrent.futures.ThreadPoolExecutor(max_workers=5) as executor:
        futures = []
        for i in range(0, len(original_texts), batch_size):
            batch = original_texts[i:i + batch_size]
            futures.append(executor.submit(translate_text_batch, batch, target_language, model_name))
            
        for future in futures:
            translated_texts.extend(future.result())
    
    # Safely apply translations back
    for i, paragraph in enumerate(paragraphs_to_translate):
        if i < len(translated_texts):
            translated_text = translated_texts[i]
            
            # Save the formatting of the first run (if any exist)
            if paragraph.runs:
                first_run = paragraph.runs[0]
                font_name = first_run.font.name
                font_size = first_run.font.size
                font_bold = first_run.font.bold
                font_italic = first_run.font.italic
                font_underline = first_run.font.underline
                try:
                    font_color = first_run.font.color.rgb if hasattr(first_run.font.color, 'rgb') else None
                except:
                    font_color = None
            else:
                font_name = font_size = font_bold = font_italic = font_underline = font_color = None

            # Clear all runs and add a new single run with the translated text
            p = paragraph._p
            for child in p.getchildren():
                if child.tag.endswith('r'):
                    p.remove(child)
            
            new_run = paragraph.add_run()
            new_run.text = translated_text
            
            # Apply saved formatting
            if font_name is not None: new_run.font.name = font_name
            if font_size is not None: new_run.font.size = font_size
            if font_bold is not None: new_run.font.bold = font_bold
            if font_italic is not None: new_run.font.italic = font_italic
            if font_underline is not None: new_run.font.underline = font_underline
            if font_color is not None: 
                try:
                    new_run.font.color.rgb = font_color
                except:
                    pass

    # Save the modified presentation
    prs.save(output_path)
    return output_path
